import streamlit as st
import os
import copy
import json
import base64
from pptx import Presentation
from pptx.util import Inches, Pt

# --- 💡 安全引入持久化组件 ---
try:
    from streamlit_javascript import st_javascript
    JS_AVAILABLE = True
except ImportError:
    JS_AVAILABLE = False

# --- 核心：无损替换单个页面文字的函数 (强制锁定字体，防止乱码/变大) ---
def safe_replace_text(text_frame, key, value):
    for paragraph in text_frame.paragraphs:
        if key in paragraph.text:
            # 捕获原文本的字体设置
            ref_font_name = "Microsoft YaHei"
            ref_font_size = Pt(14)
            if paragraph.runs:
                ref_font_name = paragraph.runs[0].font.name or ref_font_name
                ref_font_size = paragraph.runs[0].font.size or ref_font_size
            
            paragraph.text = paragraph.text.replace(key, value)
            
            # 强制重置字体属性，防止乱码和大小异常
            for run in paragraph.runs:
                run.font.name = ref_font_name
                run.font.size = ref_font_size

# --- 核心：深度克隆幻灯片的底层函数 ---
def duplicate_slide(prs, source_slide):
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in source_slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.append(new_el)
    return new_slide

# --- 核心：多页 PPT 填空与生成发动机 ---
def build_multi_page_ppt(project_title, user, date_str, problem_list):
    template_path = "template.pptx"
    if not os.path.exists(template_path): return None
    prs = Presentation(template_path)
    source_slide = prs.slides[0]
    
    for index, prob in enumerate(problem_list):
        current_slide = source_slide if index == 0 else duplicate_slide(prs, source_slide)
            
        data = {
            "{{title}}": project_title,
            "{{user}}": user,
            "{{date}}": date_str,
            "{{type}}": prob.get("type", "普通问题"),
            "{{desc}}": prob["desc"],
            "{{solve}}": prob["solve"],
            "{{duty}}": prob["duty"],
            "{{deadline}}": prob["deadline"],
            "{{decision}}": prob["decision"]
        }
        
        for shape in current_slide.shapes:
            if shape.has_text_frame:
                for key, val in data.items(): safe_replace_text(shape.text_frame, key, val)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for key, val in data.items(): safe_replace_text(cell.text_frame, key, val)
                            
        if prob.get("img_base64"):
            try:
                img_bytes = base64.b64decode(prob["img_base64"])
                temp_pic_path = f"temp_prob_{index}.jpg"
                with open(temp_pic_path, "wb") as f: f.write(img_bytes)
                current_slide.shapes.add_picture(temp_pic_path, Inches(0.52), Inches(2.3), width=Inches(2.95))
            except Exception: pass
            
    output_path = "最终汇总巡场报告.pptx"
    prs.save(output_path)
    return output_path

# --- 📱 界面逻辑 ---
st.set_page_config(page_title="设计师巡场助手", layout="centered")
st.title("📱 现场巡场助理")

if "problem_list" not in st.session_state: st.session_state.problem_list = []

team_members = ["樊洋洋", "付长春", "李新宇", "顾宇", "王硕", "郝思仆", "张晓莉", "刘璐", "吕山", "王凤国", "夏友强", "✍️ 手动输入..."]

# 1. 公共信息
project_title = st.text_input("项目名称", value="独立路壹号项目")
selected_user = st.selectbox("检查人", options=team_members)
final_user = st.text_input("请输入实际检查人") if selected_user == "✍️ 手动输入..." else selected_user
check_date = st.date_input("检查时间").strftime("%Y/%m/%d")

st.divider()

# 2. 问题录入区域
current_prob_idx = len(st.session_state.problem_list)
st.subheader(f"📷 录入巡场问题 (已添加 {current_prob_idx} 条)")

# 新增分类单选
prob_type = st.radio("巡场问题分类", options=["底线", "严控事项", "设计红黑条", "图纸错漏碰缺", "落地效果问题"], horizontal=True)

uploaded_file = st.file_uploader("拍摄照片", type=["jpg", "png"], key=f"file_{current_prob_idx}")
desc = st.text_area("问题描述", key=f"desc_{current_prob_idx}")
solve = st.text_area("解决措施", key=f"solve_{current_prob_idx}")
selected_duty = st.selectbox("责任人", options=team_members, index=7, key=f"duty_{current_prob_idx}")
final_duty = st.text_input("请输入责任人", key=f"duty_input_{current_prob_idx}") if selected_duty == "✍️ 手动输入..." else selected_duty
decision_choice = st.radio("整改决定", options=["整改", "不整改"], horizontal=True, key=f"dec_{current_prob_idx}")
deadline_date = st.date_input("完成时间", key=f"date_{current_prob_idx}").strftime("%Y/%m/%d")

if st.button("➕ 确认并添加"):
    img_b64 = base64.b64encode(uploaded_file.getbuffer()).decode("utf-8") if uploaded_file else ""
    st.session_state.problem_list.append({
        "type": prob_type, "img_base64": img_b64, "desc": desc, "solve": solve, 
        "duty": final_duty, "decision": f"{decision_choice} √", "deadline": deadline_date
    })
    st.rerun()

# 3. 汇总生成
if st.session_state.problem_list:
    if st.button("🚀 一键打包生成 PPT"):
        out = build_multi_page_ppt(project_title, final_user, check_date, st.session_state.problem_list)
        if out:
            with open(out, "rb") as f: st.download_button("📥 下载 PPT", f, file_name="报告.pptx")
        else: st.error("未找到 template.pptx 模板文件！")